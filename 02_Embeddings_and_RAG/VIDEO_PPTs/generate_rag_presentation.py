import os
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    def add_title_slide(title_text):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text

    def add_text_slide(title_text, content_text_list, font_size=Pt(20)):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        for line in content_text_list:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = font_size
        return slide

    def add_image_slide(title_text, image_path, description):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text
        left = Inches(1)
        top = Inches(1.5)
        slide.shapes.add_picture(image_path, left, top, width=Inches(7))
        txBox = slide.shapes.add_textbox(left, Inches(6.5), Inches(7), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = description
        p.font.size = Pt(18)

    # Slides
    add_title_slide("Pythonic RAG Assignment Presentation")

    overview_content = [
        "• Summarize learning points clearly.",
        "• Show concise core code snippets.",
        "• Explicit questions and answers included.",
        "• Diagrams included."
    ]
    add_text_slide("Purpose & Overview", overview_content)

    # Placeholder image slides (replace with actual image paths when running locally)
    add_image_slide("RAG Flow Diagram", "SimpleRAG.png", "Query embedding, retrieval, context, generation.")
    add_image_slide("Document Chunking Flow", "RAG_flow.png", "Document chunking to embeddings to vector DB.")

    add_text_slide("Answer #1", [
        "Q1: Can we modify embedding dimensions?",
        "A1: Yes, but only to 512 with `dimensions` parameter."
    ])

    add_text_slide("Answer #2", [
        "Q2: Benefits of async approach?",
        "A2: Faster, non-blocking, scalable, efficient for I/O tasks."
    ])

    add_text_slide("Answer #3", [
        "Q3: How to achieve reproducibility?",
        "A3: Static snapshot, temperature=0, seed parameter, iterate prompt."
    ])

    add_text_slide("Answer #4", [
        "Q4: Prompting strategies?",
        "A4: Chain of Thought, self-ask, simplify, specify format, act as expert, few-shot."
    ])

    add_text_slide("Code: PDF Loading", [
        "documents = load_pdf_documents(pdf_path + pdf_filename)",
        "if not documents: print('No documents loaded')"
    ], font_size=Pt(18))

    add_text_slide("Code: Chunking", [
        "split_documents = text_splitter.split_texts(documents)",
        "print(split_documents[0])"
    ], font_size=Pt(18))

    add_text_slide("Code: Fix Reversed Text", [
        "def fix_reversed_text(text):",
        "    corrections = {'5202': '2025', 'viXra': 'arXiv'}",
        "    for wrong, correct in corrections.items():",
        "        text = text.replace(wrong, correct)",
        "    return text"
    ], font_size=Pt(16))

    add_text_slide("Code: Metadata", [
        "metadata = {'source': 'PDF File'}",
        "vector_db.add_texts(split_documents, metadatas=[metadata]*len(split_documents))"
    ], font_size=Pt(18))

    add_text_slide("Code: Prompt Design", [
        "prompt_template = (",
        "    'Answer using context below.\\n'",
        "    'Context: {context}\\nQuestion: {question}\\nAnswer:'",
        ")"
    ], font_size=Pt(16))

    add_text_slide("Code: Distance Metrics", [
        "def euclidean_distance(vec1, vec2):",
        "    return np.sqrt(np.sum((vec1 - vec2) ** 2))",
        "def cosine_similarity(vec1, vec2):",
        "    return np.dot(vec1, vec2) / (np.linalg.norm(vec1) * np.linalg.norm(vec2))"
    ], font_size=Pt(16))

    add_text_slide("Conclusion & Learnings", [
        "- Understood RAG flow and vector DB.",
        "- Improved prompt design and metadata usage.",
        "- Implemented distance metrics.",
        "- Gained PDF handling experience."
    ])

    prs.save("Pythonic_RAG_Assignment_Presentation_Backup.pptx")

if __name__ == "__main__":
    create_presentation()

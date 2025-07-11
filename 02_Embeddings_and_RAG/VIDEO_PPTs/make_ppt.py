from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Helper function to add a slide with title and content
def add_slide(title_text, content_text_list):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = title_text
    tf = content.text_frame
    tf.clear()
    for line in content_text_list:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
    return slide

# Slide 1 — Title
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Pythonic RAG Assignment Presentation"

# Slide 2 — Purpose & Overview
content2 = [
    "• Summarize my RAG assignment and what I learned.",
    "• Describe the images and questions clearly.",
    "• Explain Activity #1 and code challenges.",
    "• Note: Corrected PDF letters read left to right manually."
]
add_slide("Purpose & Overview", content2)

# Slide 3 — Images & Flow Diagram (Placeholder text)
content3 = [
    "• Visual diagram of RAG flow.",
    "• Show document and vector database construction.",
    "• Images follow the order from the notebook."
]
add_slide("Images & Flow Diagram", content3)

# Slide 4 — Questions & Answers (Part 1)
content4 = [
    "Q1: What is RAG? How is it different from traditional QA?",
    "• RAG combines retrieval and generation to ground answers.",
    "• Traditional QA might rely only on parametric knowledge.",
    "Q2: Difference between 'query' and 'question'?",
    "• Query: a search statement, can be keywords.",
    "• Question: a full natural language sentence."
]
add_slide("Questions & Answers (1/2)", content4)

# Slide 5 — Questions & Answers (Part 2)
content5 = [
    "Q3: What are advantages of metadata?",
    "• Provides context to improve retrieval precision.",
    "• Enables filtering and organizing documents.",
    "Q4: How to include metadata in a vector DB?",
    "• Method 1: After vector DB creation (initial approach).",
    "• Method 2: Before vector DB creation (preferred)."
]
add_slide("Questions & Answers (2/2)", content5)

# Slide 6 — Activity #1 Overview
content6 = [
    "• Activity #1: Modify RAG pipeline.",
    "• Challenge 1: Adjust chunking & splitting documents.",
    "• Challenge 2: Add metadata properly.",
    "• Challenge 3: Improve prompt design and analysis."
]
add_slide("Activity #1 Overview", content6)

# Slide 7 — Code & Challenges (Part 1)
content7 = [
    "• Used CharacterTextSplitter to chunk documents.",
    "• Loaded and split texts carefully.",
    "• Printed samples to verify chunk integrity.",
    "• Handled PDF letters manually for correctness."
]
add_slide("Code & Challenges (1/2)", content7)

# Slide 8 — Code & Metadata (Part 2)
content8 = [
    "• Implemented metadata before creating vector DB.",
    "• Preferred for better filtering and relevance.",
    "• Used async embedding calls efficiently.",
    "• Designed prompts for clarity and accuracy."
]
add_slide("Code & Metadata (2/2)", content8)

# Save presentation
pptx_path = "/mnt/data/Pythonic_RAG_Assignment_Presentation_Final.pptx"
prs.save(pptx_path)

pptx_path


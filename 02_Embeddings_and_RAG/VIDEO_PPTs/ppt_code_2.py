from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create new presentation
prs = Presentation()

def add_title_slide(title_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text

def add_text_slide(title_text, content_text_list, font_size=Pt(20)):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    for line in content_text_list:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
    return slide

def add_image_slide(title_text, image_path, description):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    left = Inches(1)
    top = Inches(1.5)
    slide.shapes.add_picture(image_path, left, top, width=Inches(7))
    txBox = slide.shapes.add_textbox(left, Inches(6.5), Inches(7), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = description
    p.font.size = Pt(18)

def add_code_slide(title_text, question_text, answer_text, code_text):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p_question = tf.add_paragraph()
    p_question.text = question_text
    p_question.font.size = Pt(18)
    p_question.font.bold = True

    p_answer = tf.add_paragraph()
    p_answer.text = answer_text
    p_answer.font.size = Pt(18)

    p_code = tf.add_paragraph()
    p_code.text = code_text
    p_code.font.size = Pt(16)
    p_code.font.name = 'Courier New'

import os
import openai
# from getpass import getpass
from dotenv import load_dotenv

load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")

# Get the directory of this Python script
py_file_dir = os.path.dirname(__file__)
print(py_file_dir)

# Slide 1 — Title
add_title_slide("Pythonic RAG Assignment Presentation")

# Slide 2 — Purpose & Overview
overview_content = [
    "• Summarize my learning on RAG and vector DB concepts.",
    "• Discuss questions and detailed answers explicitly.",
    "• Show actual code snippets: chunking, PDF fixes, metadata, prompts.",
    "• Note: Corrected PDF letters manually due to left-to-right reading."
]
add_text_slide("Purpose & Overview", overview_content)

#mnt="$HOME/AIE7-BC/AIM_AIE7/02_Embeddings_and_RAG/VIDEO_PPTs"
# Slide 3 — Image 1
add_image_slide("RAG Flow Diagram", os.path.join(py_file_dir,"SimpleRAG.png"), "Query embedding, retrieval, context, and generation.")

# Slide 4 — Image 2
add_image_slide("Document Chunking & Vector Flow",os.path.join(py_file_dir,
"RAG_flow.png"), "Flow from raw docs to chunks to embeddings to vector DB.")

# Answers with questions
q1 = "Question #1: Can we modify the embedding dimensions in text-embedding-3-small? How?"
a1 = "The text-embedding-3-small model allows modifying dimensions with `dimensions` parameter, but only to 512."
c1 = "async def process_batch(batch):\n    embedding_response = await self.async_client.embeddings.create(\n        input=batch, model=self.embeddings_model_name, dimensions=512)\n    return embedding_response"

add_code_slide("Answer #1", q1, a1, c1)

q2 = "Question #2: Benefits of using an async approach for embeddings?"
a2 = "`async`: tasks run in background, faster, scalable, efficient; `sync`: blocks next task until finished."
c2 = "async def get_embeddings(list_of_text):\n    tasks = [process_batch(chunk) for chunk in list_of_text]\n    results = await asyncio.gather(*tasks)\n    return results"

add_code_slide("Answer #2", q2, a2, c2)

q3 = "Question #3: How to achieve reproducible outputs with OpenAI API?"
a3 = "1. Use static model snapshot. 2. Set temperature to 0. 3. Use seed. 4. Iterate on prompt."
c3 = "response = client.chat.completions.create(\n    model='gpt-4o',\n    messages=[...],\n    seed=123,\n    temperature=0\n)"

add_code_slide("Answer #3", q3, a3, c3)

q4 = "Question #4: Prompting strategies to guide LLM behavior?"
a4 = "Strategies: Chain of Thought, Self-ask, Explain simply, Specify format, Ask for analysis, Act as role, Few-shot."
c4 = "prompt = '''\nYou are an expert. Think step by step and explain clearly.\nFormat your answer as bullet points.\n'''"

add_code_slide("Answer #4", q4, a4, c4)

# Code logic slides
code_chunk = "text_splitter = CharacterTextSplitter()\nsplit_documents = text_splitter.split_texts(documents)\nprint(split_documents[0])"
add_text_slide("Code: Chunking & PDF Fix", [code_chunk], font_size=Pt(16))

code_meta = "metadata = {\"source\": \"PMarcaBlogs.txt\"}\nvectordb.add_texts(split_documents, metadatas=[metadata]*len(split_documents))"
add_text_slide("Code: Metadata Before Vector DB", [code_meta], font_size=Pt(16))

code_prompt = "prompt_template = (\n  'Answer the question based only on context below.\\n' \n  'Context: {context}\\nQuestion: {question}\\nAnswer:'\n)"
add_text_slide("Code: Prompt Design", [code_prompt], font_size=Pt(16))

# Conclusion
conclusion_text = [
    "Conclusion & Learnings",
    "- Mastered RAG and vector DB concepts.",
    "- Understood metadata importance and placement.",
    "- Improved prompt engineering skills.",
    "- Gained debugging experience (PDF fixes)."
]
add_text_slide("Conclusion & Learnings", conclusion_text)

# Save final presentation
#pptx_path = "/mnt/data/Pythonic_RAG_Assignment_Full_Presentation_Code.pptx"
pptx_path = os.path.join(py_file_dir,
"Pythonic_RAG_Assignment_Presentation_Code1.pptx")
prs.save(pptx_path)

pptx_path

